"""
remediate_pptx.py (courseforge) - PPTX ADA/Ally remediation engine.

Two-phase design so the AGENT supplies vision (image descriptions, slide titles)
with zero per-file human oversight:

  1) scan   : extract every issue Ally flags on PPTX + dump images to a workdir
              with a manifest the agent reads (it can VIEW the images).
  2) apply  : take the agent-written fixes.json and write a remediated copy:
              - alt text on every picture (descriptive, or "" + decorative flag)
              - slide titles (fill empty placeholder; else clone layout title /
                synthesize one, positioned OFF-CANVAS so visuals are unchanged
                but screen readers announce it - the Microsoft-documented
                "visually hidden title" remediation)
              - table header-row flag (first_row) on tables missing it
  3) verify : re-scan and report what remains (contrast issues are REPORT-ONLY;
              they are design judgments - surface to the instructor).

Usage:
  python remediate_pptx.py scan   deck.pptx --workdir W
  python remediate_pptx.py apply  deck.pptx --workdir W --out fixed.pptx
  python remediate_pptx.py verify fixed.pptx

fixes.json (written by the agent into workdir):
  {
    "alts":   { "<imageKey>": "concise description (<=110 chars)" | "" },   # "" = decorative
    "titles": { "<slideNumber>": "Title text" },
    "table_headers": true
  }
"""
import argparse, copy, json, os, re, sys

from pptx import Presentation
from pptx.util import Emu
from lxml import etree

sys.stdout.reconfigure(encoding="utf-8", errors="replace")

NS = {
    "a":   "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p":   "http://schemas.openxmlformats.org/presentationml/2006/main",
    "adec": "http://schemas.microsoft.com/office/drawing/2017/decorative",
}
DECOR_URI = "{C183D7F6-B498-43B3-948B-1728B52AA6E4}"
MAX_ALT = 110


def q(tag):
    pre, local = tag.split(":")
    return "{%s}%s" % (NS[pre], local)


# ---------- traversal ----------

def iter_pictures(shapes):
    for sh in shapes:
        if sh.shape_type == 6:  # GROUP
            yield from iter_pictures(sh.shapes)
        elif sh.shape_type == 13:  # PICTURE
            yield sh


def iter_tables(shapes):
    for sh in shapes:
        if sh.shape_type == 6:
            yield from iter_tables(sh.shapes)
        elif getattr(sh, "has_table", False):
            yield sh


def cnvpr_of(shape):
    el = shape._element
    for child in el.iter():
        if child.tag == q("p:cNvPr") or child.tag.endswith("}cNvPr"):
            return child
    return None


def get_alt(shape):
    c = cnvpr_of(shape)
    return (c.get("descr") or "").strip() if c is not None else ""


def is_decorative(shape):
    c = cnvpr_of(shape)
    if c is None:
        return False
    for ext in c.findall(".//" + q("a:ext")):
        if ext.get("uri") == DECOR_URI:
            for d in ext:
                if d.tag == q("adec:decorative") and d.get("val") in ("1", "true"):
                    return True
    return False


def slide_title_shape(slide):
    try:
        return slide.shapes.title
    except Exception:
        return None


def slide_text_snippet(slide, limit=160):
    parts = []
    for sh in slide.shapes:
        if getattr(sh, "has_text_frame", False):
            t = sh.text_frame.text.strip()
            if t:
                parts.append(t)
    return re.sub(r"\s+", " ", " | ".join(parts))[:limit]


FILENAME_RX = re.compile(r"\.(png|jpe?g|gif|bmp|svg|webp|tiff?)\s*$", re.I)


# ---------- contrast (report-only, explicit colors only) ----------

def _lum(rgb):
    def ch(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = rgb
    return 0.2126 * ch(r) + 0.7152 * ch(g) + 0.0722 * ch(b)


def contrast_ratio(fg, bg):
    l1, l2 = sorted((_lum(fg), _lum(bg)), reverse=True)
    return (l1 + 0.05) / (l2 + 0.05)


def check_contrast(slide, slide_no, issues):
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        bg = None
        try:
            if sh.fill.type == 1 and sh.fill.fore_color.type == 1:  # solid, RGB
                c = sh.fill.fore_color.rgb
                bg = (c[0], c[1], c[2])
        except Exception:
            pass
        if bg is None:
            continue  # theme/inherited fills: cannot resolve reliably - skip honestly
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                try:
                    if run.font.color and run.font.color.type == 1:
                        fc = run.font.color.rgb
                        ratio = contrast_ratio((fc[0], fc[1], fc[2]), bg)
                        if ratio < 4.5 and run.text.strip():
                            issues.append({
                                "type": "low-contrast (report-only)",
                                "slide": slide_no,
                                "detail": "text %r ratio %.2f:1 (< 4.5)" % (run.text[:30], ratio),
                            })
                except Exception:
                    pass


# ---------- scan ----------

def scan(path, workdir):
    prs = Presentation(path)
    os.makedirs(os.path.join(workdir, "images"), exist_ok=True)
    issues, images, untitled = [], [], []
    for i, slide in enumerate(prs.slides, start=1):
        for pic in iter_pictures(slide.shapes):
            alt = get_alt(pic)
            key = "s%d_id%d" % (i, pic.shape_id)
            problem = None
            if is_decorative(pic):
                pass  # explicitly decorative = fine
            elif not alt:
                problem = "image missing alt"
            elif FILENAME_RX.search(alt) or (" " not in alt and "." in alt):
                problem = "alt is a filename: %r" % alt[:40]
            elif len(alt) > MAX_ALT:
                problem = "alt too long (%d chars)" % len(alt)
            ext = pic.image.ext
            img_path = os.path.join(workdir, "images", key + "." + ext)
            with open(img_path, "wb") as f:
                f.write(pic.image.blob)
            images.append({
                "key": key, "slide": i, "path": img_path, "current_alt": alt,
                "decorative": is_decorative(pic),
                "size_px": [pic.image.size[0], pic.image.size[1]],
                "slide_context": slide_text_snippet(slide),
            })
            if problem:
                issues.append({"type": problem.split(":")[0], "slide": i,
                               "detail": problem, "key": key})
        ts = slide_title_shape(slide)
        if ts is None or not ts.text_frame.text.strip():
            issues.append({"type": "missing slide title", "slide": i,
                           "detail": "no title placeholder text"})
            untitled.append({"slide": i, "existing_text": slide_text_snippet(slide),
                             "has_empty_placeholder": ts is not None})
        for tf in iter_tables(slide.shapes):
            if not tf.table.first_row:
                issues.append({"type": "table missing header row", "slide": i,
                               "detail": "first_row flag not set"})
        check_contrast(slide, i, issues)
    report = {"deck": os.path.basename(path), "slides": len(prs.slides),
              "issues": issues, "images": images, "untitled": untitled}
    rp = os.path.join(workdir, "report.json")
    with open(rp, "w", encoding="utf-8") as f:
        json.dump(report, f, indent=1)
    hard = [x for x in issues if "report-only" not in x["type"]]
    print("SCAN %s: %d slides, %d hard issues, %d report-only, %d images extracted"
          % (report["deck"], report["slides"], len(hard),
             len(issues) - len(hard), len(images)))
    for x in issues:
        print("  [slide %2d] %s - %s" % (x["slide"], x["type"], x["detail"]))
    print("report: %s" % rp)
    return 0


# ---------- apply ----------

def set_alt(shape, text):
    c = cnvpr_of(shape)
    if c is not None:
        c.set("descr", text)


def set_decorative(shape):
    c = cnvpr_of(shape)
    if c is None:
        return
    c.set("descr", "")
    ext_lst = c.find(q("a:extLst"))
    if ext_lst is None:
        ext_lst = etree.SubElement(c, q("a:extLst"))
    for ext in ext_lst.findall(q("a:ext")):
        if ext.get("uri") == DECOR_URI:
            return
    ext = etree.SubElement(ext_lst, q("a:ext"))
    ext.set("uri", DECOR_URI)
    d = etree.SubElement(ext, q("adec:decorative"))
    d.set("val", "1")


def ensure_title(slide, prs, text):
    """Fill the empty title placeholder, or clone the layout's, or synthesize one.
    Cloned/synthesized titles are positioned OFF-CANVAS (visuals unchanged,
    screen readers announce them)."""
    ts = slide_title_shape(slide)
    if ts is not None:
        ts.text_frame.text = text
        return "filled placeholder"
    spTree = slide.shapes._spTree
    # try cloning the layout's title placeholder so it is a REAL title ph
    layout_title = None
    try:
        layout_title = slide.slide_layout.shapes.title
    except Exception:
        pass
    if layout_title is not None:
        sp = copy.deepcopy(layout_title._element)
        spTree.append(sp)
        slide_w = prs.slide_width
        xml = (
            '<a:xfrm xmlns:a="%s"><a:off x="0" y="-%d"/>'
            '<a:ext cx="%d" cy="%d"/></a:xfrm>'
            % (NS["a"], Emu(914400), slide_w, Emu(457200)))
        spPr = sp.find(".//" + q("p:spPr"))
        if spPr is None:
            spPr = sp.find(".//" + q("a:spPr"))
        if spPr is not None:
            old = spPr.find(q("a:xfrm"))
            if old is not None:
                spPr.remove(old)
            spPr.insert(0, etree.fromstring(xml))
        new_shape = slide.shapes[-1]
        new_shape.text_frame.text = text
        return "cloned layout title (off-canvas)"
    # synthesize a minimal title placeholder sp
    sp_xml = (
        '<p:sp xmlns:p="%s" xmlns:a="%s">'
        '<p:nvSpPr><p:cNvPr id="0" name="Title"/><p:cNvSpPr/>'
        '<p:nvPr><p:ph type="title"/></p:nvPr></p:nvSpPr>'
        '<p:spPr><a:xfrm><a:off x="0" y="-914400"/>'
        '<a:ext cx="%d" cy="457200"/></a:xfrm></p:spPr>'
        '<p:txBody><a:bodyPr/><a:p><a:r><a:t>%s</a:t></a:r></a:p></p:txBody>'
        '</p:sp>' % (NS["p"], NS["a"], prs.slide_width,
                     text.replace("&", "&amp;").replace("<", "&lt;")))
    spTree.append(etree.fromstring(sp_xml))
    return "synthesized title (off-canvas)"


def apply_fixes(path, workdir, out):
    fx_path = os.path.join(workdir, "fixes.json")
    with open(fx_path, encoding="utf-8") as f:
        fixes = json.load(f)
    alts = fixes.get("alts", {})
    titles = {str(k): v for k, v in fixes.get("titles", {}).items()}
    fix_tables = fixes.get("table_headers", True)
    prs = Presentation(path)
    n_alt = n_dec = n_title = n_tbl = 0
    for i, slide in enumerate(prs.slides, start=1):
        for pic in iter_pictures(slide.shapes):
            key = "s%d_id%d" % (i, pic.shape_id)
            if key in alts:
                val = alts[key].strip()
                if val == "":
                    set_decorative(pic); n_dec += 1
                else:
                    if len(val) > MAX_ALT:
                        val = val[:MAX_ALT - 1].rstrip() + "."
                    set_alt(pic, val); n_alt += 1
        if str(i) in titles:
            how = ensure_title(slide, prs, titles[str(i)])
            n_title += 1
            print("  slide %d title: %s" % (i, how))
        if fix_tables:
            for tf in iter_tables(slide.shapes):
                if not tf.table.first_row:
                    tf.table.first_row = True; n_tbl += 1
    prs.save(out)
    print("APPLY: %d alts, %d decorative, %d titles, %d table headers -> %s"
          % (n_alt, n_dec, n_title, n_tbl, out))
    return 0


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("cmd", choices=["scan", "apply", "verify"])
    ap.add_argument("deck")
    ap.add_argument("--workdir", default=None)
    ap.add_argument("--out", default=None)
    a = ap.parse_args()
    if a.cmd == "scan":
        return scan(a.deck, a.workdir or a.deck + ".work")
    if a.cmd == "apply":
        if not a.out:
            ap.error("--out required for apply")
        return apply_fixes(a.deck, a.workdir or a.deck + ".work", a.out)
    if a.cmd == "verify":
        import tempfile
        return scan(a.deck, a.workdir or tempfile.mkdtemp(prefix="pptx_verify_"))


if __name__ == "__main__":
    sys.exit(main())
